from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BLUE = RGBColor(0x1A, 0x3F, 0x6F)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xFA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x2C, 0x6D, 0xAD)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]

def add_rect(slide, l, t, w, h, color):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, text, l, t, w, h, bold=False, size=18, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = color
    return txBox

def add_bullet_box(slide, items, l, t, w, h, size=14, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = color

# ─── SLIDE 1: Titulná strana ───────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 7.5, BLUE)
add_rect(s, 0, 4.8, 13.33, 2.7, ACCENT)

add_text_box(s, "Systém na zaznamenávanie príchodov\na odchodov zamestnancov bezkontaktne",
             0.8, 1.2, 11.7, 2.8, bold=True, size=34, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "ZSI QR — Attendance Tracker",
             0.8, 4.0, 11.7, 0.7, size=20, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(s, "Martin Oros  |  FEI TUKE  |  KUI  |  2026",
             0.8, 5.3, 11.7, 0.7, size=16, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(s, "Vedúci: doc. Ing. Vojtech Čierny, CSc.",
             0.8, 6.0, 11.7, 0.6, size=14, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

# ─── SLIDE 2: Cieľ práce ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 1.1, BLUE)
add_text_box(s, "Cieľ práce", 0.4, 0.15, 12, 0.8, bold=True, size=26, color=WHITE)

add_text_box(s, "Problém", 0.4, 1.3, 5.8, 0.5, bold=True, size=18, color=BLUE)
add_bullet_box(s, [
    "• Papierové/manuálne dochádzky sú náchylné na chyby",
    "• Časovo náročná správa a archivácia",
    "• Chýba okamžitý prehľad o prítomnosti zamestnancov",
], 0.4, 1.85, 5.8, 2.2)

add_text_box(s, "Navrhnuté riešenie", 7.0, 1.3, 5.8, 0.5, bold=True, size=18, color=BLUE)
add_bullet_box(s, [
    "• Webová aplikácia pre bezkontaktné snímanie QR kódu",
    "• Automatický zápis príchodu a odchodu s časovou pečiatkou",
    "• Administrátorský prehľad, filtrovanie a export do Excelu",
    "• Generovanie unikátnych QR kódov pre každého zamestnanca",
], 7.0, 1.85, 5.8, 2.5)

add_rect(s, 0.4, 4.4, 12.5, 1.8, LIGHT_BLUE)
add_text_box(s, "Výsledok: funkčný systém, ktorý nahradí manuálne evidencie "
             "a umožní administrátorovi sledovať dochádzku v reálnom čase.",
             0.6, 4.55, 12.1, 1.4, size=15, color=DARK)

# ─── SLIDE 3: Architektúra a technológie ───────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 1.1, BLUE)
add_text_box(s, "Architektúra a technológie", 0.4, 0.15, 12, 0.8, bold=True, size=26, color=WHITE)

layers = [
    ("Klient (prehliadač)", "HTML / CSS / JavaScript\njsQR.js — dekódovanie QR kódu z kamery", 0.4),
    ("Server (PHP)", "PHP 8.0+ — spracovanie požiadaviek, autentifikácia,\nbiznis logika, generovanie QR kódov (chillerlan/php-qrcode)\nexport do XLSX", 4.6),
    ("Databáza", "MySQL / MariaDB\nTabuľky: users, employees, attendance_records, sessions", 8.8),
]
for title, desc, left in layers:
    add_rect(s, left, 1.3, 4.0, 0.6, ACCENT)
    add_text_box(s, title, left+0.1, 1.35, 3.8, 0.5, bold=True, size=14, color=WHITE)
    add_rect(s, left, 1.9, 4.0, 2.2, LIGHT_BLUE)
    add_text_box(s, desc, left+0.15, 2.0, 3.7, 2.0, size=13, color=DARK)

add_text_box(s, "Tok dát pri skenovaní:", 0.4, 4.3, 12, 0.5, bold=True, size=16, color=BLUE)
add_bullet_box(s, [
    "1.  Zamestnanec nasmeruje kameru na QR kód  →  jsQR.js dekóduje employeeId",
    "2.  POST požiadavka na PHP API /attendance/scan",
    "3.  Server overí zamestnanca a podľa posledného záznamu zapíše PRÍCHOD alebo ODCHOD",
    "4.  UI zobrazí potvrdenie s typom záznamu a časom",
], 0.4, 4.85, 12.5, 2.0, size=13)

# ─── SLIDE 4: Funkcie systému ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 1.1, BLUE)
add_text_box(s, "Funkcie systému", 0.4, 0.15, 12, 0.8, bold=True, size=26, color=WHITE)

add_text_box(s, "Zamestnanec", 0.4, 1.2, 5.8, 0.55, bold=True, size=18, color=ACCENT)
emp_items = [
    "• Prihlásenie menom a heslom",
    "• Skenovanie QR kódu kamerou zariadenia",
    "• Automatický zápis príchodu / odchodu",
    "• Zobrazenie histórie vlastnej dochádzky",
    "• Odhlásenie (zrušenie relácie)",
]
add_bullet_box(s, emp_items, 0.4, 1.8, 5.8, 3.5, size=14)

add_text_box(s, "Administrátor", 7.1, 1.2, 5.8, 0.55, bold=True, size=18, color=ACCENT)
adm_items = [
    "• Prihlásenie do oddeleného admin rozhrania",
    "• Prehľad a filtrovanie dochádzky (meno / dátum / typ)",
    "• Úprava a mazanie záznamov",
    "• Export do XLSX (Excel)",
    "• Generovanie QR kódov pre zamestnancov",
    "• Správa zamestnancov (pridanie / deaktivácia)",
]
add_bullet_box(s, adm_items, 7.1, 1.8, 5.8, 3.5, size=14)

add_rect(s, 0.4, 5.5, 12.5, 1.6, LIGHT_BLUE)
add_text_box(s, "MoSCoW prioritizácia: Must Have — prihlásenie, skenovanie QR, generovanie QR kódov, "
             "prehľad dochádzky   |   Should Have — úprava/mazanie, správa zamestnancov   |   "
             "Could Have — export do Excelu",
             0.6, 5.6, 12.1, 1.3, size=13, color=DARK)

# ─── SLIDE 5: Analýza požiadaviek (Use Case) ───────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 1.1, BLUE)
add_text_box(s, "Analýza požiadaviek — Use Case diagram", 0.4, 0.15, 12, 0.8, bold=True, size=26, color=WHITE)

add_text_box(s,
    "Systém definuje dvoch aktérov: Zamestnanca a Administrátora.\n"
    "Celkovo bolo identifikovaných 11 prípadov použitia (UC-01 – UC-11).",
    0.4, 1.2, 12.5, 0.9, size=14, color=DARK)

cols = [
    ("Zamestnanec", ["UC-01  Prihlásenie", "UC-02  Skenovanie QR kódu", "UC-03  Zobrazenie potvrdenia", "UC-04  Odhlásenie"], 0.4),
    ("Administrátor", ["UC-05  Prihlásenie admina", "UC-06  Prehľad a filtrovanie", "UC-07  Úprava záznamu", "UC-08  Vymazanie záznamu", "UC-09  Export do Excelu", "UC-10  Generovanie QR kódu", "UC-11  Správa zamestnancov"], 6.8),
]
for title, items, left in cols:
    add_rect(s, left, 2.2, 6.0, 0.55, ACCENT)
    add_text_box(s, title, left+0.15, 2.25, 5.8, 0.45, bold=True, size=15, color=WHITE)
    add_rect(s, left, 2.75, 6.0, len(items)*0.42+0.2, LIGHT_BLUE)
    add_bullet_box(s, ["  " + i for i in items], left+0.2, 2.85, 5.7, len(items)*0.42, size=13)

add_text_box(s, "Vzťahy medzi UC: «include» UC-02 → UC-03  |  «extend» UC-06 ← UC-07, UC-09",
             0.4, 6.75, 12.5, 0.55, size=13, color=ACCENT)

# ─── SLIDE 6: Záver a výsledky ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
add_rect(s, 0, 0, 13.33, 1.1, BLUE)
add_text_box(s, "Záver a dosiahnuté výsledky", 0.4, 0.15, 12, 0.8, bold=True, size=26, color=WHITE)

results = [
    ("Implementácia",
     "Funkčná webová aplikácia so zamestnaneckou a administrátorskou časťou;\n"
     "bezkontaktné skenovanie QR kódu, zápis príchodu/odchodu, prehľad dochádzky."),
    ("Testovanie",
     "15 testovacích scenárov pokrývajúcich všetky prípady použitia;\n"
     "funkčné aj negatívne scenáre — prihlásenie, sken, export, správa."),
    ("Prínos",
     "Náhrada manuálnych evidencií, okamžitý prehľad v reálnom čase,\n"
     "export dát pre mzdové účtovníctvo, jednoduché nasadenie (PHP + MySQL)."),
]
for i, (title, desc) in enumerate(results):
    top = 1.4 + i * 1.65
    add_rect(s, 0.4, top, 12.5, 0.5, ACCENT)
    add_text_box(s, title, 0.55, top+0.05, 12.0, 0.4, bold=True, size=15, color=WHITE)
    add_rect(s, 0.4, top+0.5, 12.5, 1.0, LIGHT_BLUE)
    add_text_box(s, desc, 0.55, top+0.55, 12.1, 0.9, size=13, color=DARK)

add_rect(s, 0.4, 6.4, 12.5, 0.75, BLUE)
add_text_box(s, "Ďakujem za pozornosť.  Otázky?",
             0.4, 6.45, 12.5, 0.65, bold=True, size=20, color=WHITE, align=PP_ALIGN.CENTER)

out = "/Users/martinoros/zsi-qr/employee-attendance-tracker/prezentacia_ZSI.pptx"
prs.save(out)
print("Saved:", out)
